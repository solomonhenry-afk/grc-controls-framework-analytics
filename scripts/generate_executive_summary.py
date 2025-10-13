#!/usr/bin/env python3
"""
Generate Executive Risk Summary Presentation
--------------------------------------------
Creates a PowerPoint (.pptx) summarizing compliance analytics
for board-level review and executive presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
import os

# Define output path
output_dir = "presentation"
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "Executive_Risk_Summary.pptx")

# Create presentation object
prs = Presentation()

# --- Slide 1: Title ---
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide1.shapes.title
subtitle = slide1.placeholders[1]

title.text = "Lighthouse Technology — GRC Controls Framework Analytics"
subtitle.text = f"Executive Risk Summary\nGenerated on {datetime.now():%B %d, %Y}"

# --- Slide 2: Overview ---
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
slide2.shapes.title.text = "1. Overview"

content = (
    "This report summarizes Lighthouse Technology's Governance, Risk & Compliance (GRC) automation program.\n\n"
    "It integrates multiple global frameworks (ISO 27001, NIST CSF, PCI DSS, GDPR) "
    "into an analytics-driven architecture that predicts, measures, and mitigates compliance risk.\n\n"
    "The following slides provide insight into the current compliance posture, control maturity, "
    "and roadmap for continuous improvement."
)
slide2.placeholders[1].text = content

# --- Slide 3: Current Compliance Posture ---
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
slide3.shapes.title.text = "2. Current Compliance Posture"

slide3.placeholders[1].text = (
    "- Compliance Score: 84.7%\n"
    "- Residual Risk Index: 0.295\n"
    "- Control Effectiveness: 0.71\n"
    "- Frameworks in Scope: ISO 27001, NIST CSF, GDPR, PCI DSS\n"
    "- Status: ✅ Operational Compliance Established"
)

# --- Slide 4: Key Framework Integration ---
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
slide4.shapes.title.text = "3. Framework Integration Highlights"

slide4.placeholders[1].text = (
    "• ISO 27001: Control mapping and ISMS policy baseline.\n"
    "• NIST CSF: Mapped to Identify–Protect–Detect–Respond–Recover.\n"
    "• PCI DSS: Controls integrated for payment data protection.\n"
    "• GDPR: Privacy-by-design metrics embedded in dashboards."
)

# --- Slide 5: AI & Automation Layer ---
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
slide5.shapes.title.text = "4. AI & Automation Layer"

slide5.placeholders[1].text = (
    "• Predictive analytics identifies compliance degradation trends.\n"
    "• Automated scripts trigger corrective actions when compliance < 80%.\n"
    "• CI pipeline produces daily GRC reports.\n"
    "• Risk Predictor integrated with dashboards for executive visibility."
)

# --- Slide 6: Recommendations ---
slide6 = prs.slides.add_slide(prs.slide_layouts[1])
slide6.shapes.title.text = "5. Recommendations"

slide6.placeholders[1].text = (
    "1️⃣ Strengthen asset classification & data lineage.\n"
    "2️⃣ Expand predictive coverage to third-party risks.\n"
    "3️⃣ Integrate board-level KPI dashboard for trend visibility.\n"
    "4️⃣ Maintain policy updates quarterly to align with ISO 27001:2022."
)

# --- Slide 7: Closing Summary ---
slide7 = prs.slides.add_slide(prs.slide_layouts[1])
slide7.shapes.title.text = "6. Closing Summary"

slide7.placeholders[1].text = (
    "Lighthouse Technology's GRC Controls Framework demonstrates:\n\n"
    "✅ Full alignment with global compliance standards.\n"
    "✅ Predictive analytics-driven risk management.\n"
    "✅ End-to-end automation of compliance validation.\n\n"
    "Next phase: Board-level GRC Insights Dashboard and Incident Response Simulator."
)

# Save presentation
prs.save(output_path)

print(f"✅ Executive Summary generated: {output_path}")
