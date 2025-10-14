#!/usr/bin/env python3
"""
generate_executive_summary_branding.py
Creates a fully branded executive PowerPoint with Lighthouse Technology theme.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import os

# --- Configuration ---
LOGO_PATH = "/home/solomon_henry/Downloads/ChatGPT Image Oct 14, 2025, 03_32_56 AM.png"
OUTPUT_FILE = "presentation/Executive_Risk_Summary_v2.pptx"
TAGLINE = "Empowering Smart, Predictive Governance with Automated Risk Intelligence and Resilient Cloud Compliance"
FOOTER_TEXT = "Bassey Solomon Henry | Lighthouse Technology"
BRAND_BLUE = RGBColor(10, 45, 110)
GOLD = RGBColor(212, 175, 55)
WHITE = RGBColor(255, 255, 255)

prs = Presentation()

# Utility: Add text box with formatted text
def add_textbox(slide, left, top, width, height, text, font_size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox

# Utility: Add branded background color
def set_background_color(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

# Utility: Add logo watermark (bottom-right, semi-transparent)
def add_logo_watermark(slide):
    if not os.path.exists(LOGO_PATH):
        print(f"⚠️ Logo not found at {LOGO_PATH}")
        return
    logo = slide.shapes.add_picture(LOGO_PATH, Inches(8.0), Inches(6.0), width=Inches(1.5))
    try:
        # Set transparency (semi-transparent)
        logo.fill.solid()
        logo.fill.transparency = 0.7
    except Exception:
        pass  # older pptx libs may not support transparency

# Utility: Footer text
def add_footer(slide):
    left = Inches(0.5)
    top = Inches(7.0)
    width = Inches(9)
    height = Inches(0.3)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = FOOTER_TEXT
    p.font.size = Pt(12)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.RIGHT

# --- Slide Builder ---
def add_slide(title, content, notes=None):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_background_color(slide, BRAND_BLUE)

    add_textbox(slide, Inches(0.7), Inches(0.6), Inches(9), Inches(1),
                title, font_size=32, bold=True, color=GOLD)

    if isinstance(content, list):
        top = Inches(1.8)
        for item in content:
            add_textbox(slide, Inches(1), top, Inches(8.5), Inches(0.5),
                        f"• {item}", font_size=20, color=WHITE)
            top += Inches(0.6)
    else:
        add_textbox(slide, Inches(1), Inches(2), Inches(8.5), Inches(4),
                    content, font_size=22, color=WHITE)

    add_logo_watermark(slide)
    add_footer(slide)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

# --- Slides ---

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background_color(slide, BRAND_BLUE)
add_textbox(slide, Inches(1), Inches(2), Inches(8), Inches(1.5),
            "GRC Controls Framework Analytics", font_size=40, bold=True, color=GOLD)
add_textbox(slide, Inches(1), Inches(3.3), Inches(8), Inches(1.5),
            TAGLINE, font_size=20, bold=False, color=WHITE)
add_logo_watermark(slide)
add_footer(slide)

# Business Problem
add_slide(
    "Business Problem",
    [
        "Traditional GRC systems struggle with real-time visibility.",
        "Manual audits delay executive reporting and compliance readiness.",
        "Inconsistent frameworks create fragmented risk posture awareness.",
    ],
    notes="Highlight pain points of legacy GRC approaches and lack of automation."
)

# Solution Overview
add_slide(
    "Solution Overview (STAR Framework)",
    [
        "Situation: Fragmented compliance tools limited insight.",
        "Task: Integrate automation for predictive GRC analytics.",
        "Action: Built AI-driven dashboards and automated frameworks.",
        "Result: Achieved 75% faster audit readiness and 40% reduction in manual reporting overhead.",
    ],
    notes="Walk through the STAR method for business impact demonstration."
)

# Business Impact Metrics
add_slide(
    "Business Impact Metrics",
    [
        "↑ 75% faster audit readiness",
        "↓ 40% manual compliance reporting time",
        "↑ 30% improved control maturity accuracy",
        "↗ 99.3% uptime across cloud risk workflows",
    ],
    notes="Focus on quantifiable impact metrics achieved through automation."
)

# Startup & Innovation Alignment
add_slide(
    "Startup & Innovation Focus",
    [
        "Built on Lighthouse Technology architecture for rapid scaling.",
        "Supports predictive analytics for enterprise GRC resilience.",
        "Integrates with multi-cloud, DevSecOps, and Active Directory systems.",
    ],
    notes="Position this as startup-aligned and innovation-driven architecture."
)

# Future Roadmap
add_slide(
    "Future Roadmap",
    [
        "Add Monte-Carlo Incident Response Simulator.",
        "Expand Board-Level GRC Insights Dashboard.",
        "Integrate ISO 27001 Clause-Automation Mapping Engine.",
    ],
    notes="Show readiness for iterative scaling into full GRC AI ecosystems."
)

# Save Presentation
prs.save(OUTPUT_FILE)
print(f"✅ Branded Executive Summary saved: {OUTPUT_FILE}")
