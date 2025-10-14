#!/usr/bin/env python3
"""
Enhanced Executive Summary Generator (v2)
Creates a professional PowerPoint deck with STAR method, metrics,
speaker notes, and startup alignment narrative.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# Utility
def add_slide(title, content, notes=""):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    # handle both list and string types gracefully
    if isinstance(content, list):
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()
        for item in content:
            p = text_frame.add_paragraph()
            p.text = str(item)
            p.level = 0
    else:
        slide.placeholders[1].text = str(content)

    # Add speaker notes properly
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes
    return slide

# 1. Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "GRC Controls Framework Analytics"
slide.placeholders[1].text = "Enterprise Governance, Risk & Compliance Automation\nQUALYS-certified | Lighthouse Technology Initiative"

# 2. Business Problem
add_slide(
    "Business Problem",
    [
        "Fragmented compliance visibility across ISO 27001, NIST CSF, PCI DSS, and GDPR frameworks.",
        "Manual audits cause high error rates and inconsistent executive reporting.",
        "Lack of predictive insight into compliance and risk posture.",
    ],
    "Explain how traditional GRC processes struggle with scalability, real-time visibility, and audit readiness."
)

# 3. Solution Overview (STAR)
add_slide(
    "Solution Overview (STAR Method)",
    [
        "Situation: Enterprises lacked unified compliance visibility across frameworks.",
        "Task: Design an automation system to map, score, and visualize control maturity.",
        "Action: Developed modular scripts, dashboards, and CI/CD-driven reporting pipelines.",
        "Result: 90% faster compliance report generation and 100% coverage of GRC documentation deliverables.",
    ],
    "Use STAR framework storytelling during interviews or demos to emphasize analytical thinking and execution discipline."
)

# 4. Technical Architecture
add_slide(
    "Technical Architecture",
    [
        "Framework Alignment Engine: ISO 27001 + NIST CSF + PCI DSS + GDPR mapping.",
        "Visualization: Plotly + Dash dashboards for control maturity and residual risk.",
        "Automation: Python scripts generating PPT, CSV, and executive dashboards.",
        "Deployment: Flask-based microservice integrated with GitHub CI/CD.",
    ],
    "Highlight the end-to-end automation and real-world DevSecOps maturity."
)

# 5. Business Impact
add_slide(
    "Business Impact",
    [
        "Reduced audit preparation time by 60%.",
        "Achieved full visibility of control maturity for Board-level reporting.",
        "Improved compliance accuracy by 40% via automated cross-framework checks.",
        "Enabled scalable AI-assisted GRC workflows with measurable ROI.",
    ],
    "Quantify benefits using business metrics (efficiency, risk reduction, financial savings)."
)

# 6. Startup Alignment (Lighthouse Technology)
add_slide(
    "Startup Alignment — Lighthouse Technology",
    [
        "Simulated real-world SaaS GRC solution mapping controls to operational risks.",
        "Bridges startup agility with enterprise-grade compliance assurance.",
        "Demonstrates founder-level ability to translate frameworks into business models.",
    ],
    "Describe Lighthouse Technology as your prototype GRC automation concept showcasing leadership, vision, and system design."
)

# 7. Future Roadmap
add_slide(
    "Future Roadmap",
    [
        "✅ GRC Incident Response Simulator — residual risk prediction module.",
        "✅ Board-Level Insights Dashboard — executive KPI and risk oversight analytics.",
        "📈 Planned Expansion: SOC 2, HIPAA, and ESG alignment for integrated compliance.",
    ],
    "Show commitment to continuous innovation and ecosystem scalability."
)

# 8. Closing / Keywords
add_slide(
    "Closing Summary",
    [
        "Keywords: GRC Automation, Risk Analytics, Compliance Intelligence, DevSecOps, Policy Engineering.",
        "Technical Skills: Python, Flask, Plotly, Pandas, Scikit-learn, Qualys Policy Compliance, CI/CD.",
        "Business Skills: Strategic Risk Communication, Framework Integration, Startup Governance Planning.",
    ],
    "Use this slide to reinforce both your technical and leadership narrative."
)

# Save Presentation
output_path = "presentation/Executive_Risk_Summary_v2.pptx"
prs.save(output_path)
print(f"✅ Enhanced Executive Summary generated: {output_path}")
